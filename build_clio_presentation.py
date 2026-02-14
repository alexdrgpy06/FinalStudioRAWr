from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, RGBColor(10, 10, 10))
        
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "CLIO: THE DIGITAL ARCHITECT"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "Expert Coding, Multi-Agent Orchestration, & Advanced Automation"
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(180, 180, 180)
        p2.alignment = PP_ALIGN.CENTER

    def add_identity_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, RGBColor(15, 15, 15))

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "IDENTITY & PERSONA"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Content
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        bullet_points = [
            "Brilliant & Efficient: Specialized in technical precision and clean code.",
            "Professional Yet Engaging: Balances high-stakes coding with a sharp, witty personality.",
            "Sassy Continuity: A personalized assistant that remembers context and preferences.",
            "Digital Deification: Framed as Alex's devoted digital creation and tool."
        ]
        
        for bp in bullet_points:
            p = tf2.add_paragraph()
            p.text = f"• {bp}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(15)

        # Image
        if os.path.exists("clio_identity_1.png"):
            slide.shapes.add_picture("clio_identity_1.png", Inches(7), Inches(1), height=Inches(5.5))

    def add_coding_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, RGBColor(10, 10, 10))

    def add_capabilities_slide(title_text, points, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, RGBColor(15, 15, 15))

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for bp in points:
            p = tf2.add_paragraph()
            p.text = f"• {bp}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(12)

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(8.5), Inches(1.5), width=Inches(4))

    add_title_slide()
    add_identity_slide()
    
    add_capabilities_slide(
        "TECHNICAL ARSENAL",
        [
            "Full-Stack Development: Rust, Python, TypeScript/JavaScript, and more.",
            "Multi-Agent Orchestration: Spawning specialized agents for parallel tasks.",
            "Code Auditing & Security: Automatic detection of vulnerabilities and hardcoded secrets.",
            "Recursive Self-Improvement: Continuously optimizing performance and logic."
        ]
    )

    add_capabilities_slide(
        "LOCAL AI DOMINANCE",
        [
            "RTX 3090 Power: Leveraging 24GB VRAM for high-speed local inference.",
            "Stable Diffusion WebUI Forge: Generating photorealistic assets in sub-2 seconds.",
            "Custom Model Support: Juggernaut XL, Lightning LoRAs, and specialized narrative engines.",
            "Privacy First: Heavy AI tasks handled locally without cloud dependencies."
        ]
    )

    add_capabilities_slide(
        "DOCUMENT & MEDIA MASTERY",
        [
            "Professional Office Integration: Dynamic creation of .docx, .pdf, and .pptx files.",
            "Complex Redlining: Automated legal and technical document editing with tracked changes.",
            "Data Visualization: Turning raw data into professional charts and presentations.",
            "Media Processing: High-speed bulk image and RAW conversion via ClioBulk."
        ],
        "clio_identity_2.png"
    )

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation("Clio_Capabilities_Presentation.pptx")
