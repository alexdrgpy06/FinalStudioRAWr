from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Slate, Cyber Blue, Pure White
    BG_COLOR = RGBColor(18, 18, 24)
    ACCENT_COLOR = RGBColor(0, 195, 255)
    TEXT_MAIN = RGBColor(255, 255, 255)
    TEXT_DIM = RGBColor(160, 160, 170)

    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_accents(slide):
        # Top Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.fill.background()

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide)
        
        # Decorative Element
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.8), Inches(0.1), Inches(1.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = ACCENT_COLOR
        rect.line.fill.background()

        # Main Title
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(2.7), Inches(11), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "CLIO: THE DIGITAL ARCHITECT"
        p.font.size = Pt(64)
        p.font.name = 'Arial Black'
        p.font.color.rgb = TEXT_MAIN
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(1.3), Inches(4), Inches(11), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "NEXT-GENERATION AUTONOMOUS CODING & ORCHESTRATION"
        p2.font.size = Pt(22)
        p2.font.name = 'Arial'
        p2.font.color.rgb = ACCENT_COLOR
        p2.alignment = PP_ALIGN.LEFT

    def add_content_slide(title_text, points, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide)
        add_accents(slide)

        # Slide Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN

        # Content Box
        width = Inches(6.5) if image_path else Inches(11.5)
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), width, Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for bp in points:
            p = tf2.add_paragraph()
            p.text = bp
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_DIM
            p.space_before = Pt(12)
            p.level = 0

        # Image Handling with Frame
        if image_path and os.path.exists(image_path):
            # Decorative frame for image
            left, top = Inches(7.8), Inches(1.6)
            img_width, img_height = Inches(4.8), Inches(4.8)
            
            # Accent border
            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.05), top - Inches(0.05), img_width + Inches(0.1), img_height + Inches(0.1))
            border.fill.solid()
            border.fill.fore_color.rgb = ACCENT_COLOR
            border.line.fill.background()
            
            slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

    # Build sequence
    add_title_slide()
    
    add_content_slide(
        "Identity & Core Logic",
        [
            "Clio is a high-performance digital assistant optimized for Alex's workflow.",
            "Specialized in full-stack development, security audits, and system automation.",
            "Maintains a professional and efficient demeanor as the primary operational mode.",
            "Continuous state retention ensures cross-session intelligence and consistency."
        ],
        "clio_identity_1.png"
    )

    add_content_slide(
        "Technical Capabilities",
        [
            "Multi-Agent Orchestration: Decomposition of complex tasks into parallel sub-agent runs.",
            "Security Auditing: Automated detection and remediation of vulnerabilities in Python and Rust.",
            "Performance Optimization: Direct management of C++ and Rust backends for high-speed processing.",
            "Recursive Self-Improvement: Ability to analyze and optimize its own logic and toolsets."
        ]
    )

    add_content_slide(
        "Hardware & Local AI Mastery",
        [
            "RTX 3090 Integration: Full utilization of 24GB VRAM for rapid local inference.",
            "Stable Diffusion Forge: Generation of photorealistic assets in sub-2 second cycles.",
            "Local Model Management: Specialized Juggernaut XL and Lightning LoRA configurations.",
            "Privacy-First Intelligence: Processing critical data locally to eliminate cloud dependency."
        ],
        "clio_identity_2.png"
    )

    add_content_slide(
        "Professional Office Suite",
        [
            "Dynamic Document Engineering: Automated generation of complex .docx and .pdf reports.",
            "PowerPoint Construction: Programmatic slide creation with custom themes and layouts.",
            "Data Synthesis: Transforming raw dataset inputs into professional charts and visual assets.",
            "Legal & Technical Redlining: Systematic tracked changes and document comparison logic."
        ]
    )

    prs.save(output_path)

if __name__ == "__main__":
    create_presentation("Clio_Professional_Capabilities.pptx")
